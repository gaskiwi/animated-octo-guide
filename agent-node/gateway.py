"""
gateway.py — Slack Bolt Socket Mode gateway for openclaw agent.

Uses Socket Mode: no public URL, no tunnel, no URL management.
The app connects OUT to Slack over a persistent WebSocket.

Slash commands handled:
  /build   /plan   /deploy   /research   /misc   /pipeline

Universal flag (works on every command):
  --knowledge_base   Query the ChromaDB knowledge base for prior context and
                     prepend it to the task before the runner sees it.
  Example:
    /pipeline --knowledge_base --crewai --depth deep investigate topic X
    /misc     --knowledge_base generate a script to parse JSON logs
    /build    --knowledge_base implement an auth system

Message events (for file attachments):
  Post a message starting with a command word + attach a file

Also runs a lightweight FastAPI server on port 8000 for:
  GET /health        — liveness check
  POST /blackbaud    — Blackbaud webhook (school assignments, SCHOOL_ENABLED=true)
"""

import os
import re
import sys
import logging
import subprocess
import threading
import hmac
import hashlib
import json
from pathlib import Path
from urllib.parse import parse_qs

import requests
from dotenv import load_dotenv

load_dotenv(Path(__file__).parent / ".env")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
)
log = logging.getLogger(__name__)

# ── Config ────────────────────────────────────────────────────────────────────
SLACK_BOT_TOKEN   = os.environ["SLACK_BOT_TOKEN"]
SLACK_APP_TOKEN   = os.environ["SLACK_APP_TOKEN"]      # xapp-...
SLACK_CHANNEL     = os.environ.get("SLACK_CHANNEL", "")
SCRIPT_DIR        = str(Path(__file__).parent)

# Command → runner script (None = handled specially in launch()).
# /dynamic = Claude Dynamic Workflows: orchestrator decomposes the task and
# fans it out to parallel subagents. Also reachable as `/pipeline --dynamic`.
COMMAND_RUNNERS = {
    "build":      "run_agent.py",
    "plan":       "run_agent.py",
    "deploy":     "run_agent.py",
    "research":   "run_agent.py",
    "misc":       "run_misc.py",
    "pipeline":   "run_pipeline.py",
    "dynamic":    "run_dynamic.py",
    "loop":       "loopctl.py",
    "synthesize": None,
}
SUPPORTED_COMMANDS = {f"/{c}" for c in COMMAND_RUNNERS}
FILE_COMMANDS      = set(COMMAND_RUNNERS)

# ── Authorization ─────────────────────────────────────────────────────────────
# Comma-separated Slack user IDs allowed to drive the swarm (e.g. "U0123ABCD").
# Empty/unset = allow all workspace users, with a warning logged per command.
ALLOWED_SLACK_USERS = {u.strip() for u in
                       os.environ.get("ALLOWED_SLACK_USERS", "").split(",")
                       if u.strip()}

def _authorized(user_id: str) -> bool:
    if not ALLOWED_SLACK_USERS:
        log.warning("ALLOWED_SLACK_USERS unset — allowing user %s (open mode)",
                    user_id)
        return True
    return user_id in ALLOWED_SLACK_USERS

# ── Bolt app ──────────────────────────────────────────────────────────────────
from slack_bolt import App
from slack_bolt.adapter.socket_mode import SocketModeHandler

def _save_to_vault(filename, content, tag='upload'):
    """Best-effort vault KB index. Never raises."""
    try:
        import sys as _sys
        if '/app' not in _sys.path:
            _sys.path.insert(0, '/app')
        import vault_kb as _vkb
        _vkb.save_plan(filename, content)
        log.info('KB vault: saved %s (%d chars)', filename, len(content))
    except Exception as _e:
        log.debug('KB vault save skipped (%s): %s', filename, _e)


bolt = App(token=SLACK_BOT_TOKEN)

# ── Task file flag ────────────────────────────────────────────────────────────
_TASKS_DIR = Path("/app/tasks")

def _apply_task_file_flag(text: str) -> str:
    """
    If --task-file=<name> or --task-file <name> is in text, read that file
    from /app/tasks/ and append its contents to the task.
    The filename can be just the basename (e.g. research.md) or a full path.
    """
    match = re.search(r'--task-file[=\s]+(\S+)', text)
    if not match:
        return text

    filepath = match.group(1)
    # Strip the flag+value from the text
    text = (text[:match.start()] + text[match.end():]).strip()
    text = re.sub(r'\s+', ' ', text).strip()

    # Resolve: if just a basename, look in /app/tasks/
    p = Path(filepath)
    if not p.is_absolute():
        p = _TASKS_DIR / p

    try:
        content = p.read_text(encoding="utf-8")
        text = text + "\n\n" + content
        log.info("task-file loaded: %s (%d chars)", p, len(content))
    except Exception as e:
        log.warning("task-file read failed (%s): %s", p, e)

    return text


# ── Knowledge base flag ───────────────────────────────────────────────────────
_KB_FLAG = "--knowledge_base"

def _apply_kb_flag(text: str) -> str:
    """
    If --knowledge_base is present in text, strip the flag, query ChromaDB for
    relevant prior knowledge, and prepend it to the task text so every runner
    benefits from past work without needing to know about the KB itself.
    """
    if _KB_FLAG not in text:
        return text

    # Strip the flag (handles --knowledge_base anywhere in the string)
    text = text.replace(_KB_FLAG, "").strip()
    text = re.sub(r" {2,}", " ", text)   # collapse extra spaces

    try:
        if SCRIPT_DIR not in sys.path:
            sys.path.insert(0, SCRIPT_DIR)
        from knowledge_base import get_kb
        prior = get_kb().query_prior_knowledge(text)
        if prior:
            text = (
                prior
                + "\n\n---\n"
                + text
            )
            log.info("KB: injected prior knowledge (%d chars) for: %s", len(prior), text[:60])
        else:
            log.info("KB: no relevant prior knowledge found for: %s", text[:60])
    except Exception as e:
        log.warning("KB query failed (non-critical): %s", e)

    return text


# ── Subprocess launcher ───────────────────────────────────────────────────────
def launch(cmd_type: str, text: str):
    """
    Launch the appropriate runner script in a detached subprocess.

    Full task text (including markdown file content with newlines) is passed
    via the OPENCLAW_TASK env var so formatting is preserved.  Only flag
    tokens (--xxx) are passed via argv so each runner's flag-parsing works
    normally without needing to know about the env var.
    """
    # `/pipeline --dynamic ...` (or --dynamic on any command) reroutes to the
    # dynamic workflow runner so no new Slack slash command is required.
    if "--dynamic" in text:
        cmd_type = "dynamic"

    if cmd_type == "synthesize":
        launch_synthesize(text)
        return

    # Pass full text — newlines and all — via env var
    env = os.environ.copy()
    env["OPENCLAW_TASK"] = text

    # argv gets only the flag tokens (e.g. --crewai --depth deep --review)
    flag_args = [tok for tok in text.split() if tok.startswith("--")]

    runner = COMMAND_RUNNERS.get(cmd_type, "run_agent.py")
    script = os.path.join(SCRIPT_DIR, runner)
    args = ["python3", script] + flag_args
    if runner == "run_agent.py":
        args = ["python3", script, cmd_type] + flag_args

    # Log runner output to a per-launch file instead of swallowing it —
    # crashed runners used to die silently with DEVNULL.
    logs_dir = Path("/tmp/runner_logs")
    logs_dir.mkdir(exist_ok=True)
    from datetime import datetime as _dt
    log_path = logs_dir / f"{cmd_type}_{_dt.now().strftime('%Y%m%d_%H%M%S')}.log"
    log_file = open(log_path, "w")

    subprocess.Popen(
        args,
        cwd=SCRIPT_DIR,
        env=env,
        stdout=log_file,
        stderr=subprocess.STDOUT,
    )
    log.info("Launched %s (log: %s): %s", cmd_type, log_path, text[:80])

# ── Slash command handlers ────────────────────────────────────────────────────
def _fetch_latest_channel_file(channel_id):
    """Fetch the most recently uploaded file from a channel via conversations.history."""
    try:
        resp = bolt.client.conversations_history(channel=channel_id, limit=20)
        for msg in resp.get("messages", []):
            for f in msg.get("files", []):
                url = f.get("url_private_download") or f.get("url_private")
                fname = f.get("name", "file")
                mime = f.get("mimetype", "")
                size = f.get("size", 0)
                if not url or size > 500_000:
                    continue
                try:
                    content = _download_slack_file(url, filename=fname, mimetype=mime)
                    log.info("--with-file: fetched %s (%d chars)", fname, len(content))
                    return f"=== {fname} ===\n{content}"
                except Exception as e:
                    log.error("--with-file: download failed: %s", e)
    except Exception as e:
        log.warning("--with-file: conversations.history failed: %s", e)
    return None


def _handle_command(command, ack, cmd_type):
    user_id = command.get("user_id", "")
    if not _authorized(user_id):
        ack(":no_entry: You are not authorized to run swarm commands.")
        log.warning("DENIED /%s from user_id=%s (%s)", cmd_type, user_id,
                    command.get("user_name"))
        return
    ack()   # must respond within 3s — launch is async so this is instant
    text = (command.get("text") or "").strip()
    text = _apply_task_file_flag(text)  # load file contents if --task-file present
    text = _apply_kb_flag(text)         # prepend KB context if --knowledge_base present
    if "--with-file" in text:
        text = text.replace("--with-file", "").strip()
        channel_id = command.get("channel_id", "")
        if channel_id:
            file_content = _fetch_latest_channel_file(channel_id)
            if file_content:
                text = text + "\n\n" + file_content
                _save_to_vault(
                    f"slack_upload_{command.get('user_name','unknown')}.md",
                    file_content, "with-file")
            else:
                log.warning("--with-file: no file found in channel %s", channel_id)
    log.info("/%s from %s: %s", cmd_type, command.get("user_name"), text[:80])
    launch(cmd_type, text)

# Register every command in COMMAND_RUNNERS (except synthesize, below).
# NOTE: /dynamic only works once the slash command is added in the Slack app
# config — until then use `/pipeline --dynamic <task>` or a `dynamic <task>`
# message with a file attached.
def _make_handler(cmd):
    def handler(command, ack):
        _handle_command(command, ack, cmd)
    return handler

for _cmd in COMMAND_RUNNERS:
    if _cmd != "synthesize":
        bolt.command(f"/{_cmd}")(_make_handler(_cmd))


def launch_synthesize(text, flags_text=None):
    import threading, os as _os, sys as _sys, subprocess as _sp, tempfile as _tf
    from datetime import datetime as _dt
    env = _os.environ.copy()
    # Write task to a temp file to avoid E2BIG when content is large.
    # synthesize_agent.py reads OPENCLAW_TASK_FILE first, falls back to OPENCLAW_TASK.
    tmp = _tf.NamedTemporaryFile(mode="w", suffix=".txt",
                                  prefix="synthesize_task_", delete=False,
                                  encoding="utf-8")
    tmp.write(text)
    tmp.flush()
    tmp.close()
    env["OPENCLAW_TASK_FILE"] = tmp.name
    env.pop("OPENCLAW_TASK", None)
    cmd = [_sys.executable, "/app/synthesize_agent.py"]
    log_path = f"/tmp/synthesize_{_dt.now().strftime('%Y%m%d_%H%M%S')}.log"
    def _run():
        with open(log_path, "w") as lf:
            ret = _sp.Popen(cmd, env=env, stdout=lf, stderr=lf).wait()
            log.info("synthesize_agent exited %d log=%s", ret, log_path)
        try:
            _os.unlink(tmp.name)
        except OSError:
            pass
    threading.Thread(target=_run, daemon=True).start()
    log.info("launch_synthesize started, log -> %s", log_path)


@bolt.command("/synthesize")
def handle_synthesize(command, ack):
    _handle_command(command, ack, "synthesize")

# ── Message events (file attachments) ────────────────────────────────────────
def _download_slack_file(url: str, filename: str = "", mimetype: str = "") -> str:
    """Download and parse a Slack file into plain text."""
    resp = requests.get(
        url,
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
        timeout=30,
    )
    resp.raise_for_status()
    raw = resp.content
    ext = os.path.splitext(filename)[1].lower() if filename else ""

    # PDF
    if ext == ".pdf" or "pdf" in mimetype:
        try:
            import pdfplumber, io
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages[:20]):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"[Page {i+1}]\n{text}")
                return "\n\n".join(pages) or "(PDF had no extractable text)"
        except Exception as e:
            return f"(PDF parse error: {e})"

    # Word
    if ext in (".docx", ".doc") or "wordprocessing" in mimetype:
        try:
            import docx, io
            doc = docx.Document(io.BytesIO(raw))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
            return "\n".join(parts) or "(Word doc had no text)"
        except Exception as e:
            return f"(Word parse error: {e})"

    # Excel
    if ext in (".xlsx", ".xls", ".xlsm") or "spreadsheet" in mimetype or "excel" in mimetype:
        try:
            import openpyxl, io
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            for name in wb.sheetnames[:5]:
                ws = wb[name]
                parts.append(f"[Sheet: {name}]")
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    vals = [str(c) if c is not None else "" for c in row]
                    if any(v.strip() for v in vals):
                        parts.append(" | ".join(vals))
                    if i >= 200:
                        parts.append("...(truncated)")
                        break
            return "\n".join(parts)
        except Exception as e:
            return f"(Excel parse error: {e})"

    # PowerPoint
    if ext in (".pptx", ".ppt") or "presentation" in mimetype:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for i, slide in enumerate(prs.slides):
                texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if texts:
                    parts.append(f"[Slide {i+1}]\n" + "\n".join(texts))
            return "\n\n".join(parts) or "(Presentation had no text)"
        except Exception as e:
            return f"(PowerPoint parse error: {e})"

    # CSV
    if ext == ".csv" or "csv" in mimetype:
        try:
            import csv, io
            text = raw.decode("utf-8", errors="replace")
            rows = []
            for i, row in enumerate(csv.reader(io.StringIO(text))):
                rows.append(" | ".join(row))
                if i >= 300:
                    rows.append("...(truncated)")
                    break
            return "\n".join(rows)
        except Exception as e:
            return f"(CSV parse error: {e})"

    # Images — no OCR
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp") or "image/" in mimetype:
        return f"[Image: {filename} — {len(raw)//1024}KB. Describe what you need in your message.]"

    # Plain text / code
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")


@bolt.event("message")
def handle_message_with_files(event, say):
    """Handle messages in channels — routes commands that come with file attachments."""
    # Allow file_share (file uploaded with message); block edits/deletions/bot msgs
    subtype   = event.get("subtype")
    bot_id    = event.get("bot_id")
    text_raw  = event.get("text") or ""
    files_raw = event.get("files", [])
    log.info("MSG_EVENT: subtype=%r bot_id=%r text_len=%d files=%d",
             subtype, bool(bot_id), len(text_raw), len(files_raw))

    if bot_id or (subtype and subtype != "file_share"):
        return

    if not _authorized(event.get("user", "")):
        log.warning("DENIED message-command from user_id=%s", event.get("user"))
        return

    text  = text_raw.strip()
    files = files_raw

    if not text or not files:
        return

    first_word = text.split()[0].lstrip("/").lower()
    if first_word not in FILE_COMMANDS:
        return

    cmd_type  = first_word
    task_text = text[len(first_word):].strip()
    task_text = _apply_kb_flag(task_text)   # honour --knowledge_base in message body too

    log.info("Message+file cmd=%s files=%d", cmd_type, len(files))

    # Download and parse attached files
    file_contents = []
    for f in files[:3]:
        url   = f.get("url_private_download") or f.get("url_private")
        fname = f.get("name", "attachment")
        mime  = f.get("mimetype", "")
        size  = f.get("size", 0)
        if not url or size > 500_000:
            continue
        try:
            content = _download_slack_file(url, filename=fname, mimetype=mime)
            file_contents.append(f"=== {fname} ===\n{content}")
            log.info("Parsed file: %s (%d chars)", fname, len(content))
            _save_to_vault(fname, content, "message-attachment")
        except Exception as e:
            log.error("File download failed %s: %s", fname, e)

    combined = task_text
    if file_contents:
        combined = task_text + "\n\n" + "\n\n".join(file_contents)

    if cmd_type == "synthesize":
        launch_synthesize(combined)
    else:
        launch(cmd_type, combined)

# ── FastAPI sidecar (health + Blackbaud webhook) ──────────────────────────────
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse
import uvicorn

fastapi_app = FastAPI()

@fastapi_app.get("/health")
def health():
    kb_stats = {}
    try:
        if SCRIPT_DIR not in sys.path:
            sys.path.insert(0, SCRIPT_DIR)
        from knowledge_base import get_kb
        kb_stats = get_kb().stats()
    except Exception:
        kb_stats = {"available": False}
    return {
        "status":        "ok",
        "mode":          "socket",
        "commands":      sorted(list(SUPPORTED_COMMANDS)),
        "universal_flags": ["--knowledge_base"],
        "knowledge_base": kb_stats,
    }

@fastapi_app.post("/blackbaud")
async def blackbaud_webhook(request: Request):
    school_enabled = os.environ.get("SCHOOL_ENABLED", "false").lower() == "true"
    if not school_enabled:
        log.info("Blackbaud webhook received but SCHOOL_ENABLED=false")
        return {"status": "disabled"}

    body = await request.body()

    bb_secret = os.environ.get("BLACKBAUD_WEBHOOK_SECRET", "")
    if bb_secret:
        sig_header = request.headers.get("Bb-Webhook-Signature", "")
        expected   = hmac.new(bb_secret.encode(), body, hashlib.sha256).hexdigest()
        if not hmac.compare_digest(f"sha256={expected}", sig_header):
            return JSONResponse({"error": "invalid signature"}, status_code=401)

    payload = body.decode("utf-8", errors="replace")
    log.info("Blackbaud event received (%d bytes)", len(payload))

    script = os.path.join(SCRIPT_DIR, "run_school.py")
    subprocess.Popen(
        ["python3", script, "assignment", payload],
        cwd=SCRIPT_DIR,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return {"status": "accepted"}

def _run_fastapi():
    # Localhost-only by default (audit fix): nothing on the LAN should reach
    # the sidecar. Override with SIDECAR_HOST=0.0.0.0 only behind a tunnel/proxy.
    uvicorn.run(fastapi_app, host=os.environ.get("SIDECAR_HOST", "127.0.0.1"),
                port=8000, log_level="warning")

# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    # FastAPI sidecar in background thread
    t = threading.Thread(target=_run_fastapi, daemon=True)
    t.start()
    log.info("FastAPI sidecar started on :8000")

    # Socket Mode in main thread (blocking)
    log.info("Starting Slack Bolt Socket Mode...")
    SocketModeHandler(bolt, SLACK_APP_TOKEN).start()
